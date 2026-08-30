#!/usr/bin/env python3
"""Which Claude do you use? - built on the downloaded Kourse template.

The template's formatting lives in the sample slides, not in the layouts: type
sizes, the picture bullets, the label/description run split. So rather than add
blank slides and re-style them by hand, this keeps the sample slides it needs,
swaps the text inside the existing runs, deletes the rest, and reorders what is
left. Every bit of the template's styling survives untouched.

Source template: ~/Downloads/Template.pptx (26.67 x 15in, 7 layouts).
"""
import copy
from io import BytesIO
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches

SRC = Path.home() / "Downloads" / "Template.pptx"
OUT = Path(__file__).parent / "which-claude-template.pptx"

INK = RGBColor(0x00, 0x00, 0x00)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)

WORDMARK = "tyler ai"
# Avenir Next Heavy is the nearest match on this machine to the geometric heavy
# lowercase the template's own wordmark uses.
MARK_FONT = ("/System/Library/Fonts/Avenir Next.ttc", 8)

# the template's logo images, by the part they live in
LIGHT_LOGO = "image1.png"   # white wordmark, used on the dark layouts
DARK_LOGO = "image5.png"    # black wordmark, used on the white layouts
SQUARE_MARK = "LOGO-WHITE.png"  # Kourse's square mark, top right of cover/summary


def wordmark(color, height=320):
    """Render the wordmark tight to its ink, transparent around it.

    Returns the image plus how much of it sits above the baseline. Kourse's
    wordmark has no descender, so its frame height *is* its ascent; ours has the
    'y', and matching bounding boxes would shrink our letters to fit the tail in.
    Reporting the ascent lets the caller scale on that instead, so the two
    wordmarks sit on the same baseline at the same size.
    """
    font = ImageFont.truetype(MARK_FONT[0], height, index=MARK_FONT[1])
    probe = ImageDraw.Draw(Image.new("RGBA", (1, 1)))
    x0, y0, x1, y1 = probe.textbbox((0, 0), WORDMARK, font=font, anchor="ls")
    ascent, descent = -y0, y1
    im = Image.new("RGBA", (x1 - x0, ascent + descent), (0, 0, 0, 0))
    ImageDraw.Draw(im).text((-x0, ascent), WORDMARK, font=font, fill=color,
                            anchor="ls")
    return im, ascent


def png_bytes(im):
    buf = BytesIO()
    im.save(buf, format="PNG")
    return buf.getvalue()


def rebrand(prs):
    """Swap Kourse's wordmark for ours, and drop their square mark."""
    marks = {LIGHT_LOGO: wordmark((255, 255, 255, 255)),
             DARK_LOGO: wordmark((0, 0, 0, 255))}
    metrics = {}
    for part in prs.part.package.iter_parts():
        name = str(part.partname).rsplit("/", 1)[-1]
        if name in marks:
            im, ascent = marks[name]
            part._blob = png_bytes(im)
            metrics[name] = (im.width / ascent, im.height / ascent)

    # Their frame height was all ascent. Keep that as our ascent too - same top
    # edge, same letter size - and let the frame grow right and down for the
    # extra width and the descender.
    holders = list(prs.slide_masters) + list(prs.slide_layouts)
    for holder in holders:
        for sh in list(holder.shapes):
            if "PICTURE" not in str(sh.shape_type):
                continue
            if sh.name == SQUARE_MARK:
                sh._element.getparent().remove(sh._element)
                continue
            target = str(holder.part.related_part(sh._element.blip_rId).partname)
            key = target.rsplit("/", 1)[-1]
            if key in metrics:
                w_per_ascent, h_per_ascent = metrics[key]
                ascent_in = Emu(sh.height).inches
                sh.width = Inches(ascent_in * w_per_ascent)
                sh.height = Inches(ascent_in * h_per_ascent)
                sh.name = f"{WORDMARK} wordmark"  # so it reads as ours in the editor


def paragraphs(shape):
    return shape.text_frame.paragraphs


def set_runs(par, texts):
    """Write texts into the paragraph's existing runs, keeping their formatting.

    The template's list items are two runs - a Black-weight label and a Micro
    description - so the split has to be preserved, not flattened.
    """
    runs = par.runs
    for i, t in enumerate(texts):
        if i < len(runs):
            runs[i].text = t
        else:  # clone the last run's formatting rather than inventing one
            new = copy.deepcopy(runs[-1]._r)
            runs[-1]._r.addnext(new)
            runs = par.runs
            runs[i].text = t
    for extra in runs[len(texts):]:
        extra._r.getparent().remove(extra._r)


def fill(shape, items):
    """items: list of str, or list of (label, description)."""
    pars = paragraphs(shape)
    for i, item in enumerate(items):
        set_runs(pars[i], [item] if isinstance(item, str) else list(item))
    for extra in pars[len(items):]:
        extra._p.getparent().remove(extra._p)


def recolor(slide, rgb):
    """Pin every run's colour on this slide.

    The template leaves text colour unset, so it falls through to the theme -
    where tx1 is FF0000. Nothing here is meant to be red, and the LIST slides
    sit on the master's black ground, so each slide states its own colour rather
    than inheriting a broken one.
    """
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for par in sh.text_frame.paragraphs:
            for run in par.runs:
                run.font.color.rgb = rgb


def by_name(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    raise KeyError(f"{name!r} not on slide - shapes: {[s.name for s in slide.shapes]}")


def build():
    prs = Presentation(SRC)
    s = list(prs.slides)  # 1-indexed in the comments below

    # 1 - cover
    by_name(s[0], "Title").text_frame.paragraphs[0].runs[0].text = \
        "Pick the right Claude"

    # 3 - the problem (orange full-bleed title)
    by_name(s[2], "Title one").text_frame.paragraphs[0].runs[0].text = \
        "Most wasted time comes from the wrong Claude"

    # 5 - the question (orange full-bleed title)
    by_name(s[4], "Title two").text_frame.paragraphs[0].runs[0].text = \
        "Does the job touch real files on your computer?"

    # 4 - answer: no
    by_name(s[3], "Title one").text_frame.paragraphs[0].runs[0].text = \
        "If the answer is no"
    fill(by_name(s[3], "Lorem — Ipsum…"), [
        ("Web — ", "chat in your browser, nothing to install"),
        ("Desktop app — ", "the same chat, installed. Worth it for connectors"),
        ("Either one — ", "is the right tool when you are thinking, writing, deciding"),
    ])

    # 6 - answer: yes
    by_name(s[5], "Title two").text_frame.paragraphs[0].runs[0].text = \
        "If the answer is yes"
    fill(by_name(s[5], "Lorem — Ipsum…"), [
        ("Claude Code — ", "terminal, real files. Most of this course"),
        ("It reads and writes — ", "your actual project, on disk"),
        ("It runs things — ", "commands, tests, the build"),
        ("A chat window — ", "can only hand you text to paste"),
    ])

    # 8 - the four
    by_name(s[7], "Title three").text_frame.paragraphs[0].runs[0].text = \
        "The four, plainly"
    fill(by_name(s[7], "Lorem — Ipsum…"), [
        ("Web — ", "chat in your browser. Connects to your apps"),
        ("Desktop app — ", "the same chat, installed. Worth it for connectors"),
        ("Claude Code — ", "terminal, real files. Most of this course"),
        ("Cowork — ", "the newest. The video covers it better than I can"),
    ])

    # 11 - close
    by_name(s[10], "Summary").text_frame.paragraphs[0].runs[0].text = \
        "I use all four in the same day"

    # Cover, title and summary sit on a light grey photo; the list slides sit on
    # the master's black. Colour each to its own ground.
    for i in (0, 2, 4, 10):
        recolor(s[i], INK)
    for i in (3, 5, 7):
        recolor(s[i], PAPER)

    # keep these, in this order; drop everything else
    order = [1, 3, 5, 4, 6, 8, 11]
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for i, sid in enumerate(ids, 1):
        if i not in order:
            prs.part.drop_rel(sid.rId)
            lst.remove(sid)
    keep = {i: sid for i, sid in enumerate(ids, 1) if i in order}
    for i in order:
        lst.append(keep[i])

    rebrand(prs)

    prs.save(OUT)
    print(f"saved {OUT}  ({len(list(prs.slides))} slides)")


if __name__ == "__main__":
    build()
