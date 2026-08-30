#!/usr/bin/env python3
"""Which Claude do you use? - Skool lesson deck.

Electric palette only (~/social-studio/themes/electric.json): near-white ground,
near-black ink, one blue accent used to mark the thing that matters. None of the
carousel decoration - no grid, no sticky notes, no terminal chrome. Slides for
talking over, so: few words, big type, a lot of air.

Native PowerPoint shapes throughout, so every line stays editable.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).parent / "which-claude-electric.pptx"

# electric
BG = RGBColor(0xF7, 0xF8, 0xFA)
INK = RGBColor(0x0A, 0x0A, 0x0A)
ACCENT = RGBColor(0x24, 0x54, 0xF0)
MUTED = RGBColor(0x76, 0x7C, 0x88)
RULE = RGBColor(0xDF, 0xE2, 0xE9)

HEAD = "Arial Black"
BODY = "Helvetica Neue"
MONO = "Menlo"

SW, SH = 13.333, 7.5
L = 1.15
CW = SW - L * 2
TOTAL = 7


def spacing(run, pts):
    """Letterspacing - no python-pptx API for it, so set a:rPr/@spc directly."""
    run.font._rPr.set("spc", str(int(pts * 100)))


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tf


def write(tf, parts, size, font=BODY, color=INK, bold=False, align=PP_ALIGN.LEFT,
          track=0, line=None, para=None, space_after=0):
    p = para if para is not None else tf.paragraphs[0]
    p.alignment = align
    if line:
        p.line_spacing = line
    if space_after:
        p.space_after = Pt(space_after)
    if isinstance(parts, str):
        parts = [(parts, color, bold)]
    for text, c, b in parts:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.name = font
        run.font.bold = b
        run.font.color.rgb = c
        if track:
            spacing(run, track)
    return p


def hline(slide, x, y, w, color=RULE, thick=0.012):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(thick))
    ln.shadow.inherit = False
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    return ln


def new_slide(prs, n, eyebrow=None, number=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(SH))
    bg.shadow.inherit = False
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()

    if eyebrow:
        tf = textbox(slide, L, 0.92, CW, 0.30, MSO_ANCHOR.MIDDLE)
        write(tf, eyebrow.upper(), 13, font=MONO, color=ACCENT, bold=True,
              track=2.4)
        hline(slide, L, 1.42, 0.70, color=ACCENT, thick=0.028)

    if number:
        tf = textbox(slide, SW - L - 2.0, 6.62, 2.0, 0.30, MSO_ANCHOR.MIDDLE)
        write(tf, f"{n} / {TOTAL}", 12, font=MONO, color=MUTED, track=1.2,
              align=PP_ALIGN.RIGHT)
    return slide


def headline(slide, y, parts, size=44, width=None, align=PP_ALIGN.LEFT, x=L):
    """parts: [(text, accent?), ...] joined on one paragraph; \n splits lines."""
    w = width or CW
    lines = [[]]
    for text, is_accent in parts:
        for i, chunk in enumerate(text.split("\n")):
            if i:
                lines.append([])
            if chunk:
                lines[-1].append((chunk, ACCENT if is_accent else INK, False))
    lh = size / 72 * 1.16
    for i, runs in enumerate(lines):
        tf = textbox(slide, x, y + i * lh, w, lh, MSO_ANCHOR.MIDDLE)
        write(tf, runs, size, font=HEAD, align=align, track=-0.9)
    return y + len(lines) * lh


def body(slide, y, text, size=22, width=None, color=INK, align=PP_ALIGN.LEFT,
         x=L, line=1.42):
    w = width or min(CW, 9.6)
    tf = textbox(slide, x, y, w, 2.2)
    write(tf, text, size, color=color, align=align, line=line)
    return y


def rows(slide, y, items, gap=1.02, label_w=3.5, size=21):
    """name / description, separated by hairlines. The plain list."""
    for i, (name, what) in enumerate(items):
        if i:
            hline(slide, L, y - 0.16, CW)
        tf = textbox(slide, L, y, label_w, 0.44, MSO_ANCHOR.MIDDLE)
        write(tf, name, size, color=INK, bold=True)
        tf = textbox(slide, L + label_w, y, CW - label_w, 0.44, MSO_ANCHOR.MIDDLE)
        write(tf, what, size, color=MUTED)
        y += gap
    return y


def build():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    # 1 - title
    s = new_slide(prs, 1, number=False)
    y = headline(s, 2.35, [("Pick the right\n", False), ("Claude.", True)],
                 size=54)
    body(s, y + 0.35, "Five minutes here saves you a month.", size=24,
         color=MUTED)
    tf = textbox(s, L, 6.55, CW, 0.30, MSO_ANCHOR.MIDDLE)
    write(tf, "LESSON  ·  WHICH CLAUDE", 13, font=MONO, color=MUTED, track=2.4)

    # 2 - the waste
    s = new_slide(prs, 2, eyebrow="the problem")
    y = headline(s, 2.05, [("Most wasted time comes\nfrom the ", False),
                           ("wrong Claude.", True)], size=40)
    body(s, y + 0.45, "Not from prompting badly. Not from the model. From "
                      "forcing every job through one window.", size=23,
         color=MUTED)

    # 3 - the question
    s = new_slide(prs, 3, eyebrow="one question sorts it")
    headline(s, 2.55, [("Does the job touch real\nfiles on your computer?",
                        True)], size=42)

    # 4 - no
    s = new_slide(prs, 4, eyebrow="if the answer is no")
    y = headline(s, 2.00, [("Use the ", False), ("chat window.", True)], size=42)
    body(s, y + 0.35, "You are thinking, writing, deciding. A chat window is "
                      "the right tool.", size=23, color=MUTED)
    rows(s, 4.35, [("Web", "chat in your browser"),
                   ("Desktop app", "the same chat, installed")], gap=0.92)

    # 5 - yes
    s = new_slide(prs, 5, eyebrow="if the answer is yes")
    y = headline(s, 2.00, [("Use ", False), ("Claude Code.", True)], size=42)
    body(s, y + 0.35, "Files read, written, moved, run. A chat window can only "
                      "hand you text to paste.", size=23, color=MUTED)
    rows(s, 4.35, [("Reads and writes", "your actual project, on disk"),
                   ("Runs things", "commands, tests, the build")], gap=0.92)

    # 6 - the four
    s = new_slide(prs, 6, eyebrow="what each one is")
    headline(s, 1.95, [("The four, ", False), ("plainly.", True)], size=40)
    rows(s, 3.05, [
        ("Web", "Chat in your browser. Connects to your apps."),
        ("Desktop app", "The same chat, installed. Worth it for connectors."),
        ("Claude Code", "Terminal, real files. Most of this course."),
        ("Cowork", "The newest. The video covers it better than I can."),
    ], gap=0.90, label_w=3.6, size=20)

    # 7 - close
    s = new_slide(prs, 7, eyebrow="the takeaway")
    y = headline(s, 2.15, [("I use all four\nin the ", False),
                           ("same day.", True)], size=44)
    body(s, y + 0.40, "The mistake is not picking wrong once. It is forcing "
                      "every job through one.", size=23, color=MUTED)

    prs.save(OUT)
    print(f"saved {OUT}")


if __name__ == "__main__":
    build()
